"""
File: backend/app/rag/loader.py
Purpose: Universal Multi-Format Document Loader (PDF, DOCX, TXT, MD, CSV, XLSX, PPTX).
"""

import os
import logging
from typing import Any, Dict, List

import fitz  # PyMuPDF

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".txt", ".md", ".csv", ".xlsx", ".pptx"
}


class DocumentLoaderError(Exception):
    """Custom exception for all document loading errors."""
    pass


def validate_document(file_path: str, max_size_mb: int = 50) -> bool:
    """Validates that a file exists, is within size limits, and has a supported extension."""
    if not os.path.exists(file_path):
        raise DocumentLoaderError(f"File not found: {file_path}")

    size_bytes = os.path.getsize(file_path)
    size_mb = size_bytes / (1024 * 1024)
    if size_mb > max_size_mb:
        raise DocumentLoaderError(f"File too large: {size_mb:.2f} MB. Max allowed is {max_size_mb} MB.")

    if size_bytes == 0:
        raise DocumentLoaderError("File is empty.")

    ext = os.path.splitext(file_path)[1].lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise DocumentLoaderError(
            f"Unsupported file format '{ext}'. Supported formats: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    return True


def validate_pdf(file_path: str, max_size_mb: int = 50) -> bool:
    """Validates that a file is a valid PDF and within size limits."""
    if not os.path.exists(file_path):
        raise DocumentLoaderError(f"File not found: {file_path}")

    size_bytes = os.path.getsize(file_path)
    size_mb = size_bytes / (1024 * 1024)
    if size_mb > max_size_mb:
        raise DocumentLoaderError(f"File too large: {size_mb:.2f} MB. Max allowed is {max_size_mb} MB.")

    if size_bytes == 0:
        raise DocumentLoaderError("File is empty.")

    try:
        import magic
        mime_type = magic.from_file(file_path, mime=True)
        if mime_type != 'application/pdf':
            raise DocumentLoaderError(f"Unsupported file type: {mime_type}. Only application/pdf is allowed.")
    except DocumentLoaderError:
        raise
    except Exception as e:
        logger.warning(f"Magic check failed: {e}")
        if not file_path.lower().endswith(".pdf"):
            raise DocumentLoaderError("File does not appear to be a PDF.")

    return True


def apply_ocr_to_page(page: fitz.Page) -> str:
    """Applies Optical Character Recognition to a scanned PDF page."""
    try:
        import pytesseract
        from PIL import Image
        import io

        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img = Image.open(io.BytesIO(pix.tobytes("png")))
        text = pytesseract.image_to_string(img)
        return text.strip()
    except Exception as e:
        logger.warning(f"OCR failed or Tesseract is not installed: {e}")
        return ""


def extract_text_from_pdf(file_path: str) -> Dict[str, Any]:
    """Extracts text and metadata from a PDF file."""
    try:
        doc = fitz.open(file_path)
    except fitz.FileDataError as e:
        raise DocumentLoaderError(f"File is corrupted or not a valid PDF: {e}")
    except Exception as e:
        raise DocumentLoaderError(f"Failed to open PDF: {e}")

    pages_data = []
    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text("text").strip()
        is_scanned = False

        if len(text) < 50:
            image_list = page.get_images(full=True)
            if image_list:
                is_scanned = True
                ocr_text = apply_ocr_to_page(page)
                if ocr_text:
                    text = ocr_text

        pages_data.append({
            "page_number": page_num + 1,
            "text": text,
            "is_scanned": is_scanned
        })

    metadata = doc.metadata or {}
    doc.close()

    return {
        "metadata": metadata,
        "total_pages": len(pages_data),
        "pages": pages_data
    }


def extract_text_from_docx(file_path: str) -> Dict[str, Any]:
    """Extracts paragraphs and tables from a Microsoft Word (.docx) document."""
    try:
        import docx
        doc = docx.Document(file_path)
    except Exception as e:
        raise DocumentLoaderError(f"Failed to parse Word document: {e}")

    paragraphs_text = []
    for p in doc.paragraphs:
        if p.text.strip():
            paragraphs_text.append(p.text.strip())

    for table in doc.tables:
        table_rows = []
        for row in table.rows:
            row_cells = [c.text.strip() for c in row.cells]
            table_rows.append(" | ".join(row_cells))
        if table_rows:
            paragraphs_text.append("\n".join(table_rows))

    full_text = "\n\n".join(paragraphs_text)

    # Approximate 500-word logical pages
    words = full_text.split()
    words_per_page = 400
    pages_data = []

    if not words:
        pages_data.append({"page_number": 1, "text": "", "is_scanned": False})
    else:
        for i in range(0, len(words), words_per_page):
            chunk_words = words[i:i + words_per_page]
            pages_data.append({
                "page_number": (i // words_per_page) + 1,
                "text": " ".join(chunk_words),
                "is_scanned": False
            })

    return {
        "metadata": {"title": os.path.basename(file_path)},
        "total_pages": len(pages_data),
        "pages": pages_data
    }


def extract_text_from_text_file(file_path: str) -> Dict[str, Any]:
    """Extracts text from plain text (.txt) and markdown (.md) documents."""
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
    except Exception as e:
        raise DocumentLoaderError(f"Failed to read text file: {e}")

    words = content.split()
    words_per_page = 400
    pages_data = []

    if not words:
        pages_data.append({"page_number": 1, "text": "", "is_scanned": False})
    else:
        for i in range(0, len(words), words_per_page):
            chunk_words = words[i:i + words_per_page]
            pages_data.append({
                "page_number": (i // words_per_page) + 1,
                "text": " ".join(chunk_words),
                "is_scanned": False
            })

    return {
        "metadata": {"title": os.path.basename(file_path)},
        "total_pages": len(pages_data),
        "pages": pages_data
    }


def extract_text_from_tabular(file_path: str) -> Dict[str, Any]:
    """Extracts structured tables from CSV or Excel (.csv, .xlsx) files."""
    try:
        import pandas as pd
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".csv":
            df = pd.read_csv(file_path)
        else:
            df = pd.read_excel(file_path)
    except Exception as e:
        raise DocumentLoaderError(f"Failed to parse tabular file: {e}")

    # Convert dataframe into markdown formatted tables in row batches
    rows_per_page = 50
    pages_data = []
    total_rows = len(df)

    if total_rows == 0:
        pages_data.append({"page_number": 1, "text": "Empty table", "is_scanned": False})
    else:
        for i in range(0, total_rows, rows_per_page):
            chunk_df = df.iloc[i:i + rows_per_page]
            table_md = chunk_df.to_markdown(index=False)
            pages_data.append({
                "page_number": (i // rows_per_page) + 1,
                "text": table_md or "",
                "is_scanned": False
            })

    return {
        "metadata": {"title": os.path.basename(file_path), "total_rows": total_rows},
        "total_pages": len(pages_data),
        "pages": pages_data
    }


def extract_text_from_pptx(file_path: str) -> Dict[str, Any]:
    """Extracts text slide-by-slide from PowerPoint (.pptx) presentations."""
    try:
        import pptx
        prs = pptx.Presentation(file_path)
    except Exception as e:
        raise DocumentLoaderError(f"Failed to parse PowerPoint presentation: {e}")

    pages_data = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    t = paragraph.text.strip()
                    if t:
                        slide_texts.append(t)

        pages_data.append({
            "page_number": idx,  # Slide number
            "text": "\n".join(slide_texts),
            "is_scanned": False
        })

    return {
        "metadata": {"title": os.path.basename(file_path)},
        "total_pages": len(pages_data),
        "pages": pages_data
    }


def extract_text_from_document(file_path: str) -> Dict[str, Any]:
    """Universal document extractor dispatching by file extension."""
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext in [".docx", ".doc"]:
        return extract_text_from_docx(file_path)
    elif ext in [".txt", ".md"]:
        return extract_text_from_text_file(file_path)
    elif ext in [".csv", ".xlsx"]:
        return extract_text_from_tabular(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    else:
        raise DocumentLoaderError(f"Unsupported file format: {ext}")
